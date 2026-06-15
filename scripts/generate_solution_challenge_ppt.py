from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = Path('/Users/akshatagrawal/Downloads/[EXT] Solution Challenge 2026 - Prototype PPT Template.pptx')
OUTPUT = Path('/Users/akshatagrawal/Desktop/FairFlow-AI/docs/FairFlow_Solution_Challenge_2026_Submission.pptx')

TITLE_COLOR = RGBColor(31, 41, 55)
BODY_COLOR = RGBColor(51, 65, 85)
ACCENT_BLUE = RGBColor(37, 99, 235)
ACCENT_GREEN = RGBColor(5, 150, 105)
ACCENT_AMBER = RGBColor(217, 119, 6)
WHITE = RGBColor(255, 255, 255)


def reset_text(shape, text: str, size: int = 22, bold: bool = False, color: RGBColor = TITLE_COLOR) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = 'Calibri'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, text: str) -> None:
    title_shape = slide.shapes[0]
    reset_text(title_shape, text, size=30, bold=True)


def add_multiline_box(slide, left: float, top: float, width: float, height: float, lines: list[str], *,
                      title: str | None = None, title_color: RGBColor = TITLE_COLOR, body_size: int = 18,
                      with_box: bool = False, box_color: RGBColor = RGBColor(248, 250, 252)):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if with_box else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = box_color if with_box else RGBColor(255, 255, 255)
    if with_box:
        shape.line.color.rgb = RGBColor(203, 213, 225)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = title
        r.font.name = 'Calibri'
        r.font.size = Pt(21)
        r.font.bold = True
        r.font.color.rgb = title_color

    for i, line in enumerate(lines):
        p = tf.add_paragraph() if (i > 0 or title) else tf.paragraphs[0]
        p.level = 0
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = line
        r.font.name = 'Calibri'
        r.font.size = Pt(body_size)
        r.font.color.rgb = BODY_COLOR


def add_process_step(slide, left: float, top: float, width: float, height: float, label: str, detail: str, color: RGBColor):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = label
    r1.font.name = 'Calibri'
    r1.font.bold = True
    r1.font.size = Pt(16)
    r1.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = detail
    r2.font.name = 'Calibri'
    r2.font.size = Pt(12)
    r2.font.color.rgb = WHITE


def main() -> None:
    prs = Presentation(str(TEMPLATE))

    # Slide 1: Cover
    s1 = prs.slides[0]
    reset_text(s1.shapes[0], 'FairFlow AI: Continuous Fairness Pipeline', size=34, bold=True)
    reset_text(s1.shapes[1], '', size=1)
    reset_text(s1.shapes[3], '', size=1)
    add_multiline_box(
        s1, 0.35, 1.85, 9.2, 2.1,
        [
            'Google Solution Challenge 2026 Submission',
            '',
            'Privacy-first AI fairness governance platform for hiring, lending, and healthcare.',
            'C++/WASM local precheck + FastAPI fairness engine + governance memory + mitigation workflows.',
        ],
        body_size=18,
    )
    add_multiline_box(
        s1, 0.15, 4.2, 9.2, 0.55,
        ['Build with AI Track | Problem Statement: [Unbiased AI Decision] Open Innovation'],
        body_size=15,
    )

    # Slide 2: Team details
    s2 = prs.slides[1]
    reset_text(s2.shapes[2], '', size=1)
    add_multiline_box(
        s2, 0.45, 2.75, 9.0, 1.95,
        [
            'Team Details',
            '',
            'Team name: FairFlow',
            'Team leader name: Akshat Agrawal',
            'Problem Statement: [Unbiased AI Decision] Open Innovation',
        ],
        body_size=24,
    )

    # Slide 3: Brief
    s3 = prs.slides[2]
    add_title(s3, 'Brief About Our Solution')
    add_multiline_box(
        s3, 0.55, 1.45, 8.9, 3.55,
        [
            'FairFlow AI is a full-stack fairness auditing and governance platform for AI-led decisions.',
            'It helps teams detect, explain, and mitigate bias before decisions impact real people.',
            '',
            'What we built end-to-end:',
            '• Domain-agnostic auditing for Hiring, Lending, and Healthcare datasets',
            '• Local C++/WASM privacy shield (PII hashing + schema precheck in browser)',
            '• Bias metrics (DI, SPD, Equal Opportunity, Average Odds) + SHAP + counterfactuals',
            '• Mitigation workflows (reweighing / prejudice remover / equalized odds)',
            '• Governance memory, DP-ready reporting, and certificate pipeline',
        ],
        body_size=16,
    )

    # Slide 4: Opportunities
    s4 = prs.slides[3]
    add_title(s4, 'Opportunities and Differentiation')
    add_multiline_box(
        s4, 0.55, 1.45, 2.8, 3.55,
        [
            'Most tools are offline notebooks.',
            'FairFlow offers a decision-ready UI with explainability and governance built in.',
        ],
        title='How different are we?',
        with_box=True,
    )
    add_multiline_box(
        s4, 3.55, 1.45, 2.8, 3.55,
        [
            'Local privacy precheck blocks raw PII egress.',
            'Real-time fairness diagnostics guide safe actions before rollout.',
        ],
        title='How does it solve the problem?',
        with_box=True,
    )
    add_multiline_box(
        s4, 6.55, 1.45, 2.8, 3.55,
        [
            'Domain templates + schema mapping',
            'Governance memory + mitigation recommendations',
            'Compliance-friendly PDF and audit certificate',
        ],
        title='USP',
        with_box=True,
    )

    # Slide 5: Features
    s5 = prs.slides[4]
    add_title(s5, 'List of Features Offered by the Solution')
    add_multiline_box(
        s5, 0.6, 1.45, 8.8, 3.55,
        [
            '1) Domain selector + auto schema detection (Hiring / Lending / Healthcare / Custom)',
            '2) Browser-side WASM privacy shield for PII tokenization and local precheck',
            '3) Upload pipeline with clear schema validation and guided errors',
            '4) Fairness metric engine (DI, SPD, EOD, AOD) using Fairlearn + AIF360',
            '5) Candidate-level SHAP explanation and proxy-feature flags',
            '6) Counterfactual simulation to test protected-attribute sensitivity',
            '7) Mitigation center with comparative strategy analysis',
            '8) Governance intelligence view with memory-driven risk context',
            '9) Differentially private report export + immutable certificate hash',
            '10) Authenticated multi-user dashboard with audit history',
        ],
        body_size=15,
    )

    # Slide 6: Process flow
    s6 = prs.slides[5]
    add_title(s6, 'Process Flow Diagram')
    y = 2.05
    w = 1.3
    h = 1.15
    xs = [0.45, 1.95, 3.45, 4.95, 6.45, 7.95]
    labels = [
        ('Step 1', 'Domain + CSV'),
        ('Step 2', 'WASM Precheck'),
        ('Step 3', 'Bias Audit'),
        ('Step 4', 'Explainability'),
        ('Step 5', 'Mitigation'),
        ('Step 6', 'Governance + Report'),
    ]
    colors = [ACCENT_BLUE, RGBColor(59, 130, 246), ACCENT_GREEN, RGBColor(14, 116, 144), ACCENT_AMBER, RGBColor(79, 70, 229)]
    for x, (l, d), c in zip(xs, labels, colors):
        add_process_step(s6, x, y, w, h, l, d, c)
    # arrows
    for i in range(5):
        arrow = s6.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(xs[i] + w), Inches(y + 0.43), Inches(0.45), Inches(0.28))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(100, 116, 139)
        arrow.line.fill.background()

    add_multiline_box(
        s6, 0.65, 3.5, 8.2, 1.2,
        ['Every stage is logged for traceability: input schema, metric outputs, mitigation results, and governance summary.'],
        body_size=14,
        with_box=True,
        box_color=RGBColor(239, 246, 255),
    )

    # Slide 7: Wireframes
    s7 = prs.slides[6]
    add_title(s7, 'Wireframes / Mock Screens of the Proposed Solution')
    cards = [
        (0.55, 1.55, 'A. Secure Sign-in + Workspace', 'Email auth, session token, role-ready navigation'),
        (4.85, 1.55, 'B. Upload + Domain Mapping', '2-step flow with schema checks and upload status'),
        (0.55, 3.15, 'C. Dashboard + Mitigation', 'Fairness scores, strategy comparison, recommendation'),
        (4.85, 3.15, 'D. Governance Intelligence', 'Plain-language verdict and audit memory timeline'),
    ]
    for x, y0, t, d in cards:
        add_multiline_box(s7, x, y0, 4.0, 1.35, [d], title=t, with_box=True, body_size=14)

    # Slide 8: Architecture
    s8 = prs.slides[7]
    add_title(s8, 'Architecture Diagram of the Proposed Solution')
    layers = [
        (0.7, 1.45, 8.2, 0.7, 'Client Layer: React 18 / Flutter + C++ WASM Privacy Shield', RGBColor(219, 234, 254), RGBColor(30, 64, 175)),
        (0.7, 2.25, 8.2, 0.7, 'API Layer: FastAPI + JWT + Domain Templates + Validation', RGBColor(220, 252, 231), RGBColor(6, 95, 70)),
        (0.7, 3.05, 8.2, 0.7, 'AI/ML Layer: Fairlearn, AIF360, SHAP, Counterfactuals, DoWhy, LangGraph, Gemini', RGBColor(254, 243, 199), RGBColor(146, 64, 14)),
        (0.7, 3.85, 8.2, 0.7, 'Data & Trust Layer: PostgreSQL, Audit Memories, DP Reports, SHA-256 Certificates', RGBColor(237, 233, 254), RGBColor(76, 29, 149)),
    ]
    for l, t, w0, h0, txt, fill_col, txt_col in layers:
        box = s8.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w0), Inches(h0))
        box.fill.solid()
        box.fill.fore_color.rgb = fill_col
        box.line.color.rgb = RGBColor(203, 213, 225)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = txt
        r.font.name = 'Calibri'
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = txt_col

    # Slide 9: Technologies
    s9 = prs.slides[8]
    add_title(s9, 'Technologies Used in the Solution')
    table = s9.shapes.add_table(7, 3, Inches(0.6), Inches(1.55), Inches(8.8), Inches(3.2)).table
    headers = ['Layer', 'Technologies', 'Purpose']
    rows = [
        ('Frontend', 'React 18, Tailwind, Recharts, Headless UI', 'Interactive dashboards and visual fairness analytics'),
        ('Mobile', 'Flutter (prototype track)', 'Mobile workflow for submission/use-case demos'),
        ('Backend', 'FastAPI, SQLAlchemy, JWT, Docker', 'API, auth, orchestration, and deployment-ready services'),
        ('AI/ML', 'Fairlearn, AIF360, SHAP, DoWhy', 'Metric computation, explainability, causal/proxy analysis'),
        ('AI Agent', 'LangGraph + local memory records', 'Governance recommendations with historical context'),
        ('Cloud/LLM', 'Gemini 1.5 Flash, Vertex AI, Cloud Run', 'Multimodal reasoning and scalable production path'),
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
    for r_i, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            table.cell(r_i, c).text = val
    for r in range(7):
        for c in range(3):
            tf = table.cell(r, c).text_frame
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.name = 'Calibri'
                    run.font.size = Pt(13 if r else 14)
                    run.font.bold = bool(r == 0)
                    run.font.color.rgb = TITLE_COLOR if r == 0 else BODY_COLOR

    # Slide 10: Cost
    s10 = prs.slides[9]
    add_title(s10, 'Estimated Implementation Cost (MVP to Pilot)')
    cost_table = s10.shapes.add_table(7, 3, Inches(0.8), Inches(1.55), Inches(8.4), Inches(3.25)).table
    for c, h in enumerate(['Service', 'Estimated Monthly Cost', 'Notes']):
        cost_table.cell(0, c).text = h
    cost_rows = [
        ('Cloud Run', '$0 - $20', 'Scale-to-zero for prototype traffic'),
        ('Cloud SQL / PostgreSQL', '$7 - $25', 'Small instance enough for pilot data'),
        ('Gemini / Vertex AI usage', '$5 - $30', 'Depends on Q&A/report generation volume'),
        ('Storage + certificates', '$1 - $5', 'Reports and logs'),
        ('Monitoring + logs', '$0 - $10', 'Free tier for early stage'),
        ('Total (prototype to pilot)', '$13 - $90 / month', 'Cost grows with team size and audit frequency'),
    ]
    for r_i, row in enumerate(cost_rows, start=1):
        for c, val in enumerate(row):
            cost_table.cell(r_i, c).text = val
    for r in range(7):
        for c in range(3):
            tf = cost_table.cell(r, c).text_frame
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.name = 'Calibri'
                    run.font.size = Pt(13 if r else 14)
                    run.font.bold = bool(r == 0)

    # Slide 11: MVP snapshots/metrics
    s11 = prs.slides[10]
    add_title(s11, 'Snapshots of the MVP (Implemented + Tested)')
    add_multiline_box(
        s11, 0.55, 1.5, 4.25, 3.35,
        [
            'Working modules verified on localhost:',
            '• Auth: sign-up, sign-in, protected routes',
            '• Domain templates API and 2-step upload flow',
            '• Audit pipeline with candidate explanations',
            '• Mitigation route + governance screen',
            '',
            'Sample hiring audit (200 rows):',
            'DI: 0.6585 | SPD: -0.1556 | EOD: 0.0000 | AOD: 0.0000',
            'Bias flag generated and candidate-level details stored.',
        ],
        title='Execution Proof',
        with_box=True,
        body_size=14,
    )
    add_multiline_box(
        s11, 4.95, 1.5, 4.0, 3.35,
        [
            'Domain-agnostic support validated:',
            '• Hiring schema',
            '• Lending schema',
            '• Healthcare schema',
            '',
            'Sample lending audit (200 rows):',
            'DI: 0.8010 | SPD: -0.1057 | EOD: 0.0000 | AOD: 0.0000',
            '',
            'WASM precheck and backend logs confirm end-to-end flow.',
        ],
        title='Cross-domain Validation',
        with_box=True,
        body_size=14,
    )

    # Slide 12: Future roadmap
    s12 = prs.slides[11]
    add_title(s12, 'Additional Details / Future Development')
    add_multiline_box(
        s12, 0.55, 1.5, 8.9, 3.35,
        [
            'Next 6-8 weeks roadmap:',
            '1) Async audit jobs and progress tracker for heavy SHAP/mitigation runs',
            '2) Real-time bias alerts with threshold configuration per organization',
            '3) Rich governance page: causal proxy graph + memory timeline + plain-language verdict',
            '4) Better mitigation policy selection based on fairness-accuracy tradeoff',
            '5) Judge-ready Cloud Run guest mode (no manual DB setup)',
            '6) Expanded multimodal reasoning with interview audio/video safety checks',
        ],
        body_size=16,
    )

    # Slide 13: Links
    s13 = prs.slides[12]
    tf13 = s13.shapes[0].text_frame
    tf13.clear()
    lines = [
        ('Provide links to your:', 28, True, TITLE_COLOR),
        ('', 8, False, BODY_COLOR),
        ('GitHub Public Repository', 19, True, TITLE_COLOR),
        ('https://github.com/Akshat-coder2106/FairFlow-AI', 16, False, ACCENT_BLUE),
        ('', 7, False, BODY_COLOR),
        ('Demo Video Link (3 Minutes)', 19, True, TITLE_COLOR),
        ('To be attached before final submission', 16, False, BODY_COLOR),
        ('', 7, False, BODY_COLOR),
        ('MVP Link', 19, True, TITLE_COLOR),
        ('http://127.0.0.1:3000 (local) | Cloud Run staging: publishing', 16, False, BODY_COLOR),
        ('', 7, False, BODY_COLOR),
        ('Working Prototype Link', 19, True, TITLE_COLOR),
        ('Same repository includes frontend, backend, WASM, sample datasets, and docs', 16, False, BODY_COLOR),
    ]
    for i, (txt, size, bold, color) in enumerate(lines):
        p = tf13.paragraphs[0] if i == 0 else tf13.add_paragraph()
        r = p.add_run()
        r.text = txt
        r.font.name = 'Calibri'
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

    # Slide 14: Thank you
    s14 = prs.slides[13]
    add_multiline_box(
        s14, 1.2, 1.6, 7.0, 2.2,
        ['Thank You', 'FairFlow AI Team', 'Contact: akshatagrawal21april@gmail.com'],
        title='Ready for Q&A',
        with_box=True,
        body_size=24,
        box_color=RGBColor(239, 246, 255),
        title_color=ACCENT_BLUE,
    )

    # Slide 15: Appendix
    s15 = prs.slides[14]
    add_multiline_box(
        s15, 0.8, 0.9, 8.4, 3.9,
        [
            'Fairness definitions used in the engine:',
            'DI = P(y_hat=1 | unprivileged) / P(y_hat=1 | privileged)',
            'SPD = P(y_hat=1 | unprivileged) - P(y_hat=1 | privileged)',
            'Equal Opportunity Difference = TPR_unprivileged - TPR_privileged',
            'Average Odds Difference = 0.5 * [(FPR_u - FPR_p) + (TPR_u - TPR_p)]',
            '',
            'Compliance direction: EU AI Act, NIST AI RMF, DPDP/GDPR-aligned privacy-first architecture.',
        ],
        title='Appendix: Metric Definitions and Compliance Context',
        with_box=True,
        body_size=16,
        box_color=RGBColor(248, 250, 252),
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f'Created: {OUTPUT}')


if __name__ == '__main__':
    main()
